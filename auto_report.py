#!/usr/bin/env python3
"""
服务体系 VP 月度/季度汇报自动整合工具
=============================================
用法：
    python3 auto_report.py --input ./部门合集/ --mode monthly
    python3 auto_report.py --input ./部门合集/ --mode quarterly
    python3 auto_report.py --input ./部门合集/ --mode monthly --output 整合初稿.md

输入：包含各部门 pptx 文件的目录
输出：按 SOP 框架整合的 Markdown 初稿
"""

import argparse
import os
import re
import sys
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

try:
    from pptx import Presentation
except ImportError:
    print("❌ 请先安装 python-pptx: pip install python-pptx")
    sys.exit(1)


# ============================================================
# 第一部分：核心解析模块 - 读取 pptx 提取文本+表格
# ============================================================

@dataclass
class SlideContent:
    """一页 PPT 的内容"""
    slide_number: int
    texts: list = field(default_factory=list)       # 所有文本段落
    tables: list = field(default_factory=list)      # 表格内容（每个表格是行列列表）
    title: str = ""                                  # 推断的标题（第一个短文本）

    @property
    def full_text(self) -> str:
        """获取完整文本（用于关键词匹配）"""
        parts = self.texts[:]
        for table in self.tables:
            for row in table:
                parts.append(" | ".join(row))
        return "\n".join(parts)

    def to_markdown(self) -> str:
        """转为 Markdown 格式"""
        lines = []
        if self.title:
            lines.append(f"### {self.title}")
            lines.append("")

        for text in self.texts:
            if text == self.title:
                continue
            lines.append(text)

        for table in self.tables:
            if not table:
                continue
            lines.append("")
            # 表头
            lines.append("| " + " | ".join(table[0]) + " |")
            lines.append("| " + " | ".join(["---"] * len(table[0])) + " |")
            for row in table[1:]:
                # 补齐列数
                padded = row + [""] * (len(table[0]) - len(row))
                lines.append("| " + " | ".join(padded[:len(table[0])]) + " |")
            lines.append("")

        return "\n".join(lines)


@dataclass
class PPTXFile:
    """一个 pptx 文件的解析结果"""
    filename: str
    department: str = ""    # 识别出的部门
    slides: list = field(default_factory=list)  # SlideContent 列表
    total_slides: int = 0


def extract_pptx(filepath: str) -> PPTXFile:
    """解析一个 pptx 文件，提取所有页面的文本和表格内容"""
    path = Path(filepath)
    result = PPTXFile(filename=path.name)

    try:
        prs = Presentation(str(path))
    except Exception as e:
        print(f"  ⚠️  无法解析 {path.name}: {e}")
        return result

    result.total_slides = len(prs.slides)

    for i, slide in enumerate(prs.slides, 1):
        sc = SlideContent(slide_number=i)

        # 提取文本
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        sc.texts.append(text)

            # 提取表格
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    cells = []
                    for cell in row.cells:
                        txt = cell.text_frame.text.strip().replace("\n", " | ")
                        cells.append(txt)
                    table_data.append(cells)
                sc.tables.append(table_data)

        # 推断标题（第一个长度 < 60 的文本）
        for text in sc.texts:
            if len(text) < 60 and not text.startswith("官网"):
                sc.title = text
                break

        result.slides.append(sc)

    # 识别部门
    result.department = identify_department(result)
    return result


def identify_department(pptx_file: PPTXFile) -> str:
    """根据文件名和内容识别所属部门"""
    fname = pptx_file.filename.lower()

    # 文件名匹配
    dept_patterns = {
        "人员": "人员盘点",
        "在职人员": "人员盘点",
        "spd": "SPD",
        "服务管理": "SPD",
        "bd": "BD",
        "战略合作": "BD",
        "技术支持": "技术支持部",
        "ts": "技术支持部",
        "售后服务": "售后服务部",
        "售后": "售后服务部",
        "交付中心": "技术交付中心",
        "交付": "技术交付中心",
        "rayman": "架构师",
    }

    for pattern, dept in dept_patterns.items():
        if pattern in fname:
            return dept

    # 内容匹配（看前3页）
    content = ""
    for slide in pptx_file.slides[:3]:
        content += slide.full_text + "\n"

    content_lower = content.lower()
    if "技术交付中心" in content:
        return "技术交付中心"
    elif "售后服务部" in content:
        return "售后服务部"
    elif "技术支持" in content:
        return "技术支持部"
    elif "spd" in content_lower or "服务管理" in content:
        return "SPD"
    elif "bd" in content_lower or "战略合作" in content:
        return "BD"

    return "未识别"


# ============================================================
# 第二部分：内容匹配模块 - 关键词映射到 SOP 槽位
# ============================================================

@dataclass
class SlotMatch:
    """一个槽位的匹配结果"""
    slot_id: str
    slot_name: str
    department: str           # 期望来源部门
    matched_slides: list = field(default_factory=list)  # (PPTXFile, SlideContent) 列表
    is_quarterly_only: bool = False


# 月度槽位定义：(slot_id, slot_name, 期望部门, 匹配关键词列表)
MONTHLY_SLOTS = [
    # 板块一：人员盘点
    ("1-1", "组织架构图", "人员盘点", ["人员盘点", "服务体系VP"]),
    ("1-2", "在岗/入离职情况", "人员盘点", ["在岗人数", "入职", "离职", "在职情况"]),
    ("1-3", "编制 vs 在职表", "人员盘点", ["编制", "在职", "地区"]),
    ("1-4", "奋斗者计划", "人员盘点", ["奋斗者"]),

    # 板块二：战略指标 OGSM
    ("2-1", "万元服务成本", "SPD", ["万元业绩服务成本", "万元服务成本"]),
    ("2-2", "交付成本控制", "SPD|技术交付中心", ["交付成本控制", "标准杆人天", "人均回款"]),
    ("2-3", "售前失误/特批/非项目人天", "SPD|技术交付中心", ["售前失误人天", "特批人天", "非项目人天"]),
    ("2-4", "驻场客户签约", "技术交付中心", ["驻场客户", "驻场", "到期"]),
    ("2-5", "售后服务效率", "SPD|售后服务部", ["NPS", "服务满意度", "超时占比", "工单重开率"]),

    # 板块三：关键数据
    ("3-1", "在服客户数", "售后服务部", ["在服客户数"]),
    ("3-2", "非KA平均交付损益人天", "技术交付中心", ["非KA项目平均交付", "损益人天"]),
    ("3-3", "非KA项目周期", "技术交付中心", ["非KA项目.*周期", "项目周期"]),
    ("3-4", "Top20客户工时", "售后服务部", ["Top20客户", "TOP.*20", "工时.*占所有"]),
    ("3-5", "P1&S1问题统计", "售后服务部", ["P1.*S1", "P1问题", "S1.*安全问题", "账号被盗"]),

    # 板块四：AI关键项目
    ("4-0", "AI提效关键项目进度总览", "SPD", ["AI提效.*项目.*总览", "AI提效关键"]),
    ("4-1", "文档中台", "技术交付中心", ["文档中台", "DocHub", "智能方案撰写"]),
    ("4-2", "交付工具集CMBox", "技术交付中心", ["交付工具", "CMBox", "CMbox"]),
    ("4-3", "新小C同学", "SPD", ["小C同学", "小 C 同学"]),
    ("4-4", "数字人葱葱", "售后服务部", ["数字人葱葱", "葱葱"]),
    ("4-5", "双儿底座", "技术支持部", ["双儿底座", "双儿.*安全", "凭证到执行"]),
    ("4-6", "双儿使用情况", "技术支持部", ["双儿使用", "会话总数.*独立联系人", "会话总数"]),
    ("4-7", "双儿开发情况", "技术支持部", ["双儿开发", "开发计划.*已完成"]),
    ("4-8", "CMCU自动化审核", "SPD", ["CMCU自动化", "CMCU.*审核"]),
    ("4-9", "云服务中心自动开通", "SPD", ["云服务中心.*开通", "自动开通", "自动续期"]),

    # 板块五：重点工作
    ("5-1", "重点SOP/Skill项目", "技术支持部", ["Air客户端.*SOP", "SOP.*Skill", "完成Air"]),
    ("5-2", "巡检进度", "售后服务部", ["巡检进度", "巡检.*完成"]),
    ("5-3", "BD项目申报", "BD", ["项目申报.*资金", "申报.*资金申请"]),
    ("5-4", "BD政府关系与行业交流", "BD", ["政府关系.*行业交流", "行业交流"]),
    ("5-5", "BD业务拓展与客户跟进", "BD", ["业务拓展.*客户跟进", "业务拓展"]),

    # 板块六：工作计划
    ("6-1", "非KA业绩冲刺", "技术交付中心", ["非KA.*目标总盘", "非KA.*回款", "业绩冲刺"]),
    ("6-2", "售前AI计划", "技术交付中心", ["AI.*工作计划", "AI 工作计划"]),
    ("6-3", "售后-服务中台", "售后服务部", ["服务中台"]),
    ("6-4", "售后-质检", "售后服务部", ["售后质检", "质检规则", "质检"]),
    ("6-5", "售后-自动邮件通知", "售后服务部", ["自动邮件通知", "批量邮件"]),
    ("6-6", "售后-客户分级模型", "售后服务部", ["客户分级", "客户价值.*服务成本"]),
    ("6-7", "售后TS-SOP补全计划", "技术支持部", ["完整树.*SOP", "SOP补全"]),
    ("6-8", "SPD计划", "SPD", ["服务管理AI化", "服务产品管理", "OGSM指标"]),
    ("6-9", "BD计划", "BD", ["BD.*工作计划", "项目申报工作计划", "跟进沟通.*项目"]),
]

# 季度额外槽位
QUARTERLY_EXTRA_SLOTS = [
    ("Q1-1", "员工结构分析", "人员盘点", ["员工结构分析", "性别.*司龄", "平均年龄"]),
    ("Q1-2", "入离职明细汇总", "人员盘点", ["入离职.*明细", "主动离职.*被动"]),
    ("Q3-1", "二线工单季度同比", "技术支持部", ["二线.*工单.*同比", "二线工单"]),
    ("Q5-1", "HW加固", "售后服务部", ["HW.*加固", "安全加固", "HW客户"]),
    ("Q6-1", "人天管理中台化", "售后服务部", ["人天管理.*中台", "人天.*中台化"]),
    ("Q6-2", "商机挖掘", "技术支持部", ["商机挖掘", "智能Agent.*商机", "商机转化"]),
    ("Q6-3", "海外计划", "技术支持部", ["海外计划", "出海.*合规", "海外知识库"]),
    ("Q6-4", "技术中台建设", "技术支持部", ["技术中台", "排班管理.*晚班", "效率看板"]),
    ("Q4-1", "im-bot审批工作流", "SPD", ["im-bot.*审批", "报价审批"]),
    ("Q4-2", "自动回访", "SPD|售后服务部", ["自动回访", "自动触发回访"]),
]


def match_slide_to_slots(slide: SlideContent, department: str, slots: list) -> list:
    """将一页 PPT 匹配到可能的槽位"""
    matched = []
    full_text = slide.full_text

    for slot_id, slot_name, expected_dept, keywords in slots:
        # 检查部门是否匹配
        dept_match = False
        for d in expected_dept.split("|"):
            if d == department or department == "未识别":
                dept_match = True
                break
        if not dept_match:
            continue

        # 关键词匹配
        for kw in keywords:
            try:
                if re.search(kw, full_text):
                    matched.append(slot_id)
                    break
            except re.error:
                if kw in full_text:
                    matched.append(slot_id)
                    break

    return matched


def build_slot_mapping(pptx_files: list, mode: str = "monthly") -> dict:
    """
    构建槽位映射：slot_id -> [(pptx_file, slide), ...]
    """
    slots = MONTHLY_SLOTS[:]
    if mode == "quarterly":
        slots.extend(QUARTERLY_EXTRA_SLOTS)

    # 初始化映射
    mapping = {}
    for slot_id, slot_name, dept, kws in slots:
        mapping[slot_id] = SlotMatch(
            slot_id=slot_id,
            slot_name=slot_name,
            department=dept,
            is_quarterly_only=slot_id.startswith("Q")
        )

    # 遍历所有文件的所有页
    for pptx_file in pptx_files:
        for slide in pptx_file.slides:
            matched_slots = match_slide_to_slots(slide, pptx_file.department, slots)
            for slot_id in matched_slots:
                mapping[slot_id].matched_slides.append((pptx_file, slide))

    return mapping


# ============================================================
# 第三部分：整合输出模块 - 生成 Markdown 初稿
# ============================================================

def generate_report(mapping: dict, mode: str = "monthly", month: str = "") -> str:
    """按 SOP 框架生成整合初稿 Markdown"""

    lines = []

    # 封面
    if mode == "monthly":
        lines.append(f"# {month}总结 AND 下月计划")
    else:
        lines.append(f"# {month}总结 AND 下季度计划")
    lines.append("")
    lines.append("Marvin")
    lines.append("")
    lines.append("---")
    lines.append("")

    # 定义输出顺序
    sections = [
        ("## 一、人员盘点", ["1-1", "1-2", "1-3", "1-4"]),
    ]

    # 季度增量
    if mode == "quarterly":
        sections[0] = ("## 一、人员盘点", ["1-1", "Q1-1", "1-2", "Q1-2", "1-3", "1-4"])

    sections.extend([
        ("## 二、战略指标（OGSM）", ["2-1", "2-2", "2-3", "2-4", "2-5"]),
        ("## 三、关键数据", ["3-1", "3-2", "3-3", "3-4", "3-5"]),
    ])

    if mode == "quarterly":
        sections[-1] = ("## 三、关键数据", ["3-1", "3-2", "3-3", "Q3-1", "3-4", "3-5"])

    # AI 关键项目
    ai_slots = ["4-0", "4-1", "4-2", "4-3", "4-4", "4-5", "4-6", "4-7", "4-8", "4-9"]
    if mode == "quarterly":
        ai_slots.extend(["Q4-1", "Q4-2"])
    sections.append(("## 四、AI 关键项目（核心）", ai_slots))

    # 重点工作
    work_slots = ["5-1", "5-2", "5-3", "5-4", "5-5"]
    if mode == "quarterly":
        work_slots.insert(2, "Q5-1")
    sections.append(("## 五、重点工作", work_slots))

    # 工作计划
    plan_slots = ["6-1", "6-2", "6-3", "6-4", "6-5", "6-6", "6-7", "6-8", "6-9"]
    if mode == "quarterly":
        plan_slots.extend(["Q6-1", "Q6-2", "Q6-3", "Q6-4"])
    sections.append(("## 六、工作计划", plan_slots))

    # 生成各板块
    for section_title, slot_ids in sections:
        lines.append(section_title)
        lines.append("")

        for slot_id in slot_ids:
            if slot_id not in mapping:
                continue
            slot = mapping[slot_id]

            # 槽位标题
            quarterly_tag = " 🔷季度" if slot.is_quarterly_only else ""
            lines.append(f"### 📄 {slot.slot_name}{quarterly_tag}")
            lines.append(f"> 来源部门：{slot.department}")
            lines.append("")

            if not slot.matched_slides:
                lines.append("⚠️ **未匹配到内容** — 请手动补充")
                lines.append("")
            else:
                for pptx_file, slide in slot.matched_slides:
                    lines.append(f"*来源：{pptx_file.filename} P{slide.slide_number}*")
                    lines.append("")
                    lines.append(slide.to_markdown())
                    lines.append("")

            lines.append("---")
            lines.append("")

    # 统计信息
    lines.append("## 📊 整合统计")
    lines.append("")
    total_matched = sum(1 for s in mapping.values() if s.matched_slides)
    total_slots = len(mapping)
    lines.append(f"- 总槽位数：{total_slots}")
    lines.append(f"- 已匹配：{total_matched}")
    lines.append(f"- 未匹配：{total_slots - total_matched}")
    lines.append("")

    # 未匹配清单
    unmatched = [s for s in mapping.values() if not s.matched_slides]
    if unmatched:
        lines.append("### ⚠️ 未匹配槽位（需手动补充）")
        lines.append("")
        for slot in unmatched:
            lines.append(f"- [ ] {slot.slot_id} {slot.slot_name}（期望来源：{slot.department}）")
        lines.append("")

    # 未使用的页面
    lines.append("### 📋 未使用的部门 PPT 页面（已过滤）")
    lines.append("")
    all_used_slides = set()
    for slot in mapping.values():
        for pptx_file, slide in slot.matched_slides:
            all_used_slides.add((pptx_file.filename, slide.slide_number))

    for pptx_file in pptx_files_global:
        unused = []
        for slide in pptx_file.slides:
            if (pptx_file.filename, slide.slide_number) not in all_used_slides:
                title_hint = slide.title[:40] if slide.title else "(无标题)"
                unused.append(f"P{slide.slide_number}: {title_hint}")
        if unused:
            lines.append(f"**{pptx_file.filename}**（{pptx_file.department}）")
            for u in unused:
                lines.append(f"  - {u}")
            lines.append("")

    return "\n".join(lines)


# ============================================================
# 第四部分：主入口
# ============================================================

# 全局变量用于 generate_report 访问
pptx_files_global = []


def main():
    parser = argparse.ArgumentParser(
        description="服务体系 VP 月度/季度汇报自动整合工具",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python3 auto_report.py --input ./部门合集/ --mode monthly
  python3 auto_report.py --input ./部门合集/ --mode quarterly --month "Q1"
  python3 auto_report.py --input ./部门合集/ --mode monthly --month "4月" --output 4月整合初稿.md
        """
    )
    parser.add_argument("--input", "-i", required=True, help="包含部门 pptx 文件的目录路径")
    parser.add_argument("--mode", "-m", choices=["monthly", "quarterly"], default="monthly",
                        help="整合模式：monthly（月度，默认）或 quarterly（季度）")
    parser.add_argument("--month", default="本期", help="报告月份/季度标识，如 '4月' 或 'Q1'")
    parser.add_argument("--output", "-o", default=None, help="输出 Markdown 文件路径（默认打印到终端）")

    args = parser.parse_args()

    input_dir = Path(args.input)
    if not input_dir.exists():
        print(f"❌ 输入目录不存在：{input_dir}")
        sys.exit(1)

    # 查找所有 pptx 文件
    pptx_paths = sorted(input_dir.glob("*.pptx"))
    if not pptx_paths:
        # 尝试子目录
        pptx_paths = sorted(input_dir.rglob("*.pptx"))

    # 过滤掉 __MACOSX 目录
    pptx_paths = [p for p in pptx_paths if "__MACOSX" not in str(p)]

    if not pptx_paths:
        print(f"❌ 未在 {input_dir} 中找到 pptx 文件")
        sys.exit(1)

    print(f"📂 输入目录：{input_dir}")
    print(f"📋 模式：{'月度' if args.mode == 'monthly' else '季度'}")
    print(f"📄 找到 {len(pptx_paths)} 个 pptx 文件：")

    # 解析所有文件
    global pptx_files_global
    pptx_files = []
    for path in pptx_paths:
        print(f"  🔍 解析：{path.name} ...", end=" ")
        result = extract_pptx(str(path))
        print(f"✅ {result.total_slides}页 → 识别为【{result.department}】")
        pptx_files.append(result)

    pptx_files_global = pptx_files

    # 构建槽位映射
    print(f"\n🔗 构建内容映射...")
    mapping = build_slot_mapping(pptx_files, mode=args.mode)

    matched_count = sum(1 for s in mapping.values() if s.matched_slides)
    total_count = len(mapping)
    print(f"  ✅ 已匹配：{matched_count}/{total_count} 个槽位")

    # 生成报告
    print(f"\n📝 生成整合初稿...")
    report = generate_report(mapping, mode=args.mode, month=args.month)

    # 输出
    if args.output:
        output_path = Path(args.output)
        output_path.write_text(report, encoding="utf-8")
        print(f"\n✅ 已输出到：{output_path}")
        print(f"   文件大小：{output_path.stat().st_size / 1024:.1f} KB")
    else:
        print("\n" + "=" * 80)
        print(report)

    print(f"\n🎉 整合完成！共 {matched_count}/{total_count} 个槽位已填充。")
    if matched_count < total_count:
        print(f"   ⚠️  有 {total_count - matched_count} 个槽位未匹配，请检查输出中的提示手动补充。")


if __name__ == "__main__":
    main()
